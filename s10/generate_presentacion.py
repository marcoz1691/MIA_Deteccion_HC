"""Genera presentación PowerPoint del proyecto CITIMED — detección de inconsistencias."""
from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT_DIR = ROOT / "s10" / "docs"
EVID = ROOT / "s10" / "evidencias"
CAPTURAS = EVID / "capturas"

# Paleta CITIMED
COLOR_PRIMARY = RGBColor(0x0D, 0x47, 0xA1)   # azul médico
COLOR_ACCENT = RGBColor(0x00, 0x96, 0x88)    # teal
COLOR_DARK = RGBColor(0x26, 0x32, 0x38)
COLOR_MUTED = RGBColor(0x54, 0x6E, 0x7A)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT = RGBColor(0xEC, 0xF0, 0xF1)


def load_metricas() -> dict:
    resumen_path = EVID / "resumen_metricas.json"
    if resumen_path.exists():
        return json.loads(resumen_path.read_text(encoding="utf-8"))
    return {
        "tfidf_test": {
            "roc_auc": 0.9485,
            "auprc": 0.4186,
            "localizacion_top1": 0.8457,
            "notas_con_error": 311,
        }
    }


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_bar(slide, title: str, subtitle: str = "") -> None:
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.75))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.4))
        sf = sb.text_frame
        sf.clear()
        sp = sf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(14)
        sp.font.color.rgb = COLOR_MUTED


def add_bullets(slide, items: list[str], left=0.6, top=1.6, width=12.0, height=5.5, size=20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = COLOR_DARK
        p.space_after = Pt(10)


def add_two_column(slide, left_items: list[str], right_items: list[str], left_title="", right_title=""):
    if left_title:
        lt = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(5.8), Inches(0.4))
        ltf = lt.text_frame
        ltf.text = left_title
        ltf.paragraphs[0].font.bold = True
        ltf.paragraphs[0].font.size = Pt(18)
        ltf.paragraphs[0].font.color.rgb = COLOR_ACCENT

    if right_title:
        rt = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.4))
        rtf = rt.text_frame
        rtf.text = right_title
        rtf.paragraphs[0].font.bold = True
        rtf.paragraphs[0].font.size = Pt(18)
        rtf.paragraphs[0].font.color.rgb = COLOR_ACCENT

    add_bullets(slide, left_items, left=0.6, top=1.95, width=5.8, height=4.8, size=17)
    add_bullets(slide, right_items, left=6.8, top=1.95, width=5.8, height=4.8, size=17)


def add_metric_cards(slide, metrics: list[tuple[str, str, str]], top=2.0):
    card_w = 2.8
    gap = 0.35
    start_x = (13.33 - (len(metrics) * card_w + (len(metrics) - 1) * gap)) / 2
    for i, (value, label, detail) in enumerate(metrics):
        x = start_x + i * (card_w + gap)
        card = slide.shapes.add_shape(1, Inches(x), Inches(top), Inches(card_w), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_LIGHT
        card.line.color.rgb = COLOR_ACCENT

        vb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(top + 0.25), Inches(card_w - 0.3), Inches(0.8))
        vf = vb.text_frame
        vf.text = value
        vf.paragraphs[0].font.size = Pt(36)
        vf.paragraphs[0].font.bold = True
        vf.paragraphs[0].font.color.rgb = COLOR_PRIMARY
        vf.paragraphs[0].alignment = PP_ALIGN.CENTER

        lb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(top + 1.05), Inches(card_w - 0.3), Inches(0.5))
        lf = lb.text_frame
        lf.text = label
        lf.paragraphs[0].font.size = Pt(14)
        lf.paragraphs[0].font.bold = True
        lf.paragraphs[0].font.color.rgb = COLOR_DARK
        lf.paragraphs[0].alignment = PP_ALIGN.CENTER

        if detail:
            db = slide.shapes.add_textbox(Inches(x + 0.15), Inches(top + 1.55), Inches(card_w - 0.3), Inches(0.5))
            df = db.text_frame
            df.text = detail
            df.paragraphs[0].font.size = Pt(11)
            df.paragraphs[0].font.color.rgb = COLOR_MUTED
            df.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_image_if_exists(slide, path: Path, left, top, width):
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))


def build_presentation() -> Presentation:
    metricas = load_metricas()
    m = metricas.get("tfidf_test", {})
    roc = m.get("roc_auc", 0.949)
    auprc = m.get("auprc", 0.419)
    loc = m.get("localizacion_top1", 0.846)
    n_err = m.get("notas_con_error", 311)
    aciertos = int(loc * n_err)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Portada ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_PRIMARY)

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.5))
    tf = title_box.text_frame
    tf.text = "Detección de inconsistencias\nen historias clínicas"
    for p in tf.paragraphs:
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(11.7), Inches(0.8))
    sf = sub.text_frame
    sf.text = "Proyecto CITIMED · MIA Detección HC"
    sf.paragraphs[0].font.size = Pt(24)
    sf.paragraphs[0].font.color.rgb = RGBColor(0xB2, 0xEB, 0xF2)
    sf.paragraphs[0].alignment = PP_ALIGN.CENTER

    team = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.7), Inches(1.2))
    tff = team.text_frame
    tff.text = (
        "Patricio Bayas Meza · José Puebla Paladines · Marco Zurita Rojas\n"
        "2026 · github.com/marcoz1691/MIA_Deteccion_HC"
    )
    for p in tff.paragraphs:
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER

    # --- Slide 2: Problema ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "El problema", "¿Por qué importa detectar inconsistencias?")
    add_bullets(slide, [
        "Las historias clínicas pueden contener errores de documentación que afectan la seguridad del paciente.",
        "Ejemplos: medicación contraindicada, diagnóstico incompatible con el plan, datos de laboratorio incoherentes.",
        "Revisar manualmente cada nota es lento, costoso y propenso a pasar errores por alto.",
        "Objetivo: asistir al clínico señalando la oración sospechosa, no reemplazar su juicio.",
    ])

    # --- Slide 3: Qué es una inconsistencia ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "¿Qué es una inconsistencia clínica?")
    add_two_column(
        slide,
        left_title="Quién la crea",
        right_title="Quién la detecta",
        left_items=[
            "Errores del profesional al redactar la nota.",
            "Contradicciones entre secciones (antecedentes, diagnóstico, plan).",
            "Incompatibilidad con hechos médicos o guías clínicas.",
            "En evaluación: anotadores del dataset MEDEC marcan la oración errónea.",
        ],
        right_items=[
            "El sistema automático (TF-IDF + LLM + RAG).",
            "Segmenta la nota en oraciones y puntúa cada una (0–1).",
            "Devuelve top-1: la oración más sospechosa.",
            "El médico valida la alerta y decide la acción clínica.",
        ],
    )

    # --- Slide 4: Insight clave ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Insight clave", "De nivel nota → nivel oración")
    add_bullets(slide, [
        "En MEDEC, la nota con error y la corregida comparten ~96,6 % del texto idéntico.",
        "La oración errónea representa solo ~8,9 % del contenido total.",
        "Clasificar toda la nota (baseline S6): ROC-AUC ≈ 0,504 → equivalente al azar.",
        "Reformular a nivel oración (S7): ROC-AUC ≈ 0,949 y localización top-1 ≈ 84,6 %.",
        "Conclusión: detectar inconsistencias exige razonamiento semántico, no solo léxico.",
    ], top=1.55, size=19)

    # --- Slide 5: Arquitectura ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Arquitectura del sistema")
    add_bullets(slide, [
        "Entrada: texto libre de la historia clínica.",
        "1. Segmentación → oraciones numeradas (sid).",
        "2. Scoring por brazo: TF-IDF · LLM zero-shot · LLM + RAG (FAISS + guías clínicas).",
        "3. Fusión de scores → localización top-1 (65 % LLM + 35 % TF-IDF).",
        "4. Salida: JSON con scores, alerta y oración a revisar.",
        "Interfaces: API FastAPI (8000) + React (5173) · Demo Streamlit (8501).",
    ], top=1.55, size=18)

    add_image_if_exists(slide, CAPTURAS / "frontend_analisis.png", left=7.2, top=1.55, width=5.5)

    # --- Slide 6: Tres brazos ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Tres brazos de inferencia")
    add_two_column(
        slide,
        left_title="TF-IDF (S6)",
        right_title="LLM + RAG (S7)",
        left_items=[
            "Modelo léxico: n-gramas + regresión logística.",
            "Entrenado sobre oraciones MEDEC con Error Sentence ID.",
            "Rápido (~1–5 ms/oración), 100 % local.",
            "Rol: baseline y fallback si la API LLM falla.",
        ],
        right_items=[
            "LLM zero-shot: ¿esta oración es inconsistente? (SI/NO).",
            "LLM + RAG: contexto de guías (medicación, diagnóstico, odontología).",
            "Razonamiento semántico anclado en conocimiento externo.",
            "Modo demo: mock LLM local sin API key.",
        ],
    )

    # --- Slide 7: Ejemplo ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Ejemplo: error de medicación")
    add_bullets(slide, [
        "Paciente con alergia documentada a penicilina.",
        "Oración sospechosa: «Se indica amoxicilina 500 mg cada 8 h pese a alergia documentada a penicilina.»",
        "El sistema resalta esa oración con score alto.",
        "RAG recupera fragmentos de guías de farmacoterapia que justifican la alerta.",
        "El clínico revisa y corrige la prescripción.",
    ], top=1.55, width=6.5, size=18)
    add_image_if_exists(slide, CAPTURAS / "demo_analisis.png", left=7.0, top=1.55, width=5.8)

    # --- Slide 8: Resultados ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Resultados principales (MEDEC test)")
    add_metric_cards(slide, [
        (f"{roc:.3f}", "ROC-AUC", "Nivel oración (TF-IDF ajustado)"),
        (f"{auprc:.3f}", "AUPRC", f"Prevalencia ~4,5 %"),
        (f"{loc*100:.1f} %", "Localización top-1", f"{aciertos}/{n_err} notas con error"),
        ("0,504", "ROC-AUC baseline", "Nivel nota (S6, azar)"),
    ], top=2.1)
    note = slide.shapes.add_textbox(Inches(0.8), Inches(4.6), Inches(11.7), Inches(1.5))
    nf = note.text_frame
    nf.text = (
        "Reformular la tarea de nota completa a oración individual eleva el rendimiento "
        "de azar (0,504) a discriminación clínica útil (0,949)."
    )
    nf.paragraphs[0].font.size = Pt(16)
    nf.paragraphs[0].font.color.rgb = COLOR_MUTED
    nf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_image_if_exists(slide, EVID / "figura_ajuste.png", left=3.5, top=5.2, width=6.3)

    # --- Slide 9: Comparación tripartita ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Comparación tripartita", "TF-IDF vs LLM zero-shot vs LLM + RAG")
    add_bullets(slide, [
        "TF-IDF supera al LLM mock en este corpus (AUC ~0,95 vs ~0,50).",
        "El mock LLM simula respuestas; con API real se espera mejor desempeño semántico.",
        "Arquitectura híbrida recomendada: TF-IDF pre-filtra → LLM+RAG en candidatas.",
        "Fallback producción: si la API cae → TF-IDF solo + alerta al usuario.",
    ], top=1.55, width=5.5, size=17)
    add_image_if_exists(
        slide,
        CAPTURAS / "figura_comparacion_tripartita.png",
        left=6.2,
        top=1.45,
        width=6.5,
    )

    # --- Slide 10: Por tipo de error ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Análisis por tipo de error")
    add_bullets(slide, [
        "Tipos MEDEC: diagnosis, pharmacotherapy, management, causalOrganism, testResult.",
        "Mayor recall en diagnosis (~90 %) y pharmacotherapy (~86 %).",
        "Management es el más difícil (~68 % recall).",
        "Permite priorizar mejoras por dominio clínico.",
    ], top=1.55, width=5.5, size=17)
    add_image_if_exists(slide, EVID / "heatmap_recall_por_tipo.png", left=6.0, top=1.45, width=6.7)

    # --- Slide 11: Interfaces ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Prototipo funcional")
    add_two_column(
        slide,
        left_title="FastAPI + React",
        right_title="Streamlit (demo)",
        left_items=[
            "POST /generar → análisis JSON.",
            "OpenAPI en /docs.",
            "UI moderna con scores por brazo.",
            "Proxy Vite → localhost:8000.",
        ],
        right_items=[
            "Demo para presentación en clase.",
            "Ejemplos precargados (medicación, diagnóstico).",
            "Página de métricas y configuración.",
            "Consentimiento PHI al usar API real.",
        ],
    )
    add_image_if_exists(slide, CAPTURAS / "frontend_inicio.png", left=0.6, top=4.8, width=5.5)
    add_image_if_exists(slide, CAPTURAS / "demo_metricas.png", left=6.8, top=4.8, width=5.5)

    # --- Slide 12: Producción ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Riesgos y despliegue")
    add_bullets(slide, [
        "PHI (CITIMED): anonimización local + Ollama on-premise; no enviar datos a cloud sin DPA.",
        "Latencia: cascada TF-IDF → LLM reduce de ~8 s a ~2 s por nota.",
        "Costo API: ~$0,08–0,15 por 1.000 oraciones (gpt-4o-mini).",
        "Modo degradado: API caída → TF-IDF + banner de alerta (sin mock silencioso).",
        "Auditoría: hash SHA-256 en audit.log, sin texto clínico en logs.",
    ], top=1.55, size=18)
    add_image_if_exists(slide, CAPTURAS / "demo_fallback.png", left=7.5, top=1.55, width=5.2)

    # --- Slide 13: Conclusiones ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Conclusiones")
    add_bullets(slide, [
        "Detectar inconsistencias clínicas no es un problema léxico: el baseline TF-IDF a nivel nota no supera el azar.",
        "La granularidad oración + Error Sentence ID de MEDEC desbloquea métricas útiles (AUC 0,949, localización 84,6 %).",
        "Arquitectura híbrida TF-IDF + LLM + RAG equilibra velocidad, costo y razonamiento semántico.",
        "Prototipo end-to-end listo: API, frontend React, demo Streamlit y pipeline CITIMED con anonimizador.",
        "Próximo paso: piloto odontológico con Ollama local y validación clínica en casos reales anonimizados.",
    ], top=1.55, size=18)

    # --- Slide 14: Cierre ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_PRIMARY)
    thanks = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.2))
    tff = thanks.text_frame
    tff.text = "¡Gracias!"
    tff.paragraphs[0].font.size = Pt(48)
    tff.paragraphs[0].font.bold = True
    tff.paragraphs[0].font.color.rgb = COLOR_WHITE
    tff.paragraphs[0].alignment = PP_ALIGN.CENTER

    demo = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(11.7), Inches(2.0))
    dff = demo.text_frame
    dff.text = (
        "Demo en vivo:\n"
        "streamlit run demo/app.py  ·  uvicorn api.main:app --port 8000\n\n"
        "Repositorio: github.com/marcoz1691/MIA_Deteccion_HC"
    )
    for p in dff.paragraphs:
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0xB2, 0xEB, 0xF2)
        p.alignment = PP_ALIGN.CENTER

    return prs


def build_presentation_5min() -> Presentation:
    """Versión corta (~5 min): 8 diapositivas esenciales."""
    metricas = load_metricas()
    m = metricas.get("tfidf_test", {})
    roc = m.get("roc_auc", 0.949)
    auprc = m.get("auprc", 0.419)
    loc = m.get("localizacion_top1", 0.846)
    n_err = m.get("notas_con_error", 311)
    aciertos = int(loc * n_err)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # 1. Portada
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_PRIMARY)
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.5))
    tf = title_box.text_frame
    tf.text = "Detección de inconsistencias\nen historias clínicas"
    for p in tf.paragraphs:
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER
    sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(11.7), Inches(1.0))
    sf = sub.text_frame
    sf.text = "CITIMED · Bayas · Puebla · Zurita · 2026"
    sf.paragraphs[0].font.size = Pt(22)
    sf.paragraphs[0].font.color.rgb = RGBColor(0xB2, 0xEB, 0xF2)
    sf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 2. Problema
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "El problema", "~5 min · asistir al clínico, no reemplazarlo")
    add_bullets(slide, [
        "Errores en historias clínicas: medicación contraindicada, plan incompatible con diagnóstico.",
        "Revisión manual: lenta y propensa a pasar errores por alto.",
        "Objetivo: señalar la oración sospechosa para que el médico decida.",
    ], top=1.55, size=22)

    # 3. Insight clave
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Insight clave", "Nota completa → oración individual")
    add_bullets(slide, [
        "MEDEC: nota con error vs. corregida comparten ~96,6 % del texto.",
        "Baseline a nivel nota: ROC-AUC = 0,504 (azar).",
        "A nivel oración: ROC-AUC = 0,949 · localización top-1 = 84,6 %.",
        "→ No es un problema léxico; exige razonamiento semántico.",
    ], top=1.55, size=21)

    # 4. Cómo funciona
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Cómo funciona")
    add_bullets(slide, [
        "1. Segmentar nota en oraciones.",
        "2. Tres brazos: TF-IDF (rápido, local) · LLM zero-shot · LLM + RAG (guías clínicas).",
        "3. Fusión de scores → oración top-1 a revisar.",
        "API FastAPI + React · Demo Streamlit · fallback TF-IDF si cae el LLM.",
    ], top=1.55, width=6.2, size=19)
    add_image_if_exists(slide, CAPTURAS / "frontend_analisis.png", left=6.8, top=1.45, width=6.0)

    # 5. Ejemplo
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Ejemplo")
    add_bullets(slide, [
        "Alergia a penicilina documentada.",
        "Oración sospechosa: amoxicilina pese a alergia.",
        "Sistema resalta la oración · RAG justifica con guías de farmacoterapia.",
    ], top=1.55, width=6.0, size=20)
    add_image_if_exists(slide, CAPTURAS / "demo_analisis.png", left=6.5, top=1.45, width=6.2)

    # 6. Resultados
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Resultados (MEDEC test)")
    add_metric_cards(slide, [
        (f"{roc:.3f}", "ROC-AUC", "Nivel oración"),
        (f"{auprc:.3f}", "AUPRC", "Prevalencia 4,5 %"),
        (f"{loc*100:.1f} %", "Localización", f"{aciertos}/{n_err} notas"),
    ], top=2.0)
    add_image_if_exists(slide, EVID / "figura_ajuste.png", left=2.5, top=4.5, width=8.3)

    # 7. Prototipo y cierre técnico
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Prototipo y próximos pasos")
    add_bullets(slide, [
        "Prototipo end-to-end: API + React + Streamlit (mock LLM local).",
        "Anonimizador CITIMED para PHI · Ollama on-premise en producción.",
        "Próximo paso: piloto odontológico con validación clínica.",
        "github.com/marcoz1691/MIA_Deteccion_HC",
    ], top=1.55, size=20)
    add_image_if_exists(slide, CAPTURAS / "demo_metricas.png", left=7.0, top=1.55, width=5.5)

    # 8. Gracias
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_PRIMARY)
    thanks = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.5))
    tff = thanks.text_frame
    tff.text = "¡Gracias!\n¿Preguntas?"
    for i, p in enumerate(tff.paragraphs):
        p.font.size = Pt(44 if i == 0 else 28)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER

    return prs


def main():
    import argparse

    parser = argparse.ArgumentParser(description="Genera presentación PowerPoint CITIMED")
    parser.add_argument(
        "--duracion",
        choices=["completa", "5min"],
        default="completa",
        help="completa (~14 diapos, ~12 min) o 5min (8 diapos, ~5 min)",
    )
    args = parser.parse_args()

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    if args.duracion == "5min":
        out_path = OUT_DIR / "Presentacion_CITIMED_5min.pptx"
        prs = build_presentation_5min()
    else:
        out_path = OUT_DIR / "Presentacion_CITIMED.pptx"
        prs = build_presentation()

    prs.save(str(out_path))
    print(f"Presentación generada: {out_path}")
    print(f"Diapositivas: {len(prs.slides)}")


if __name__ == "__main__":
    main()
